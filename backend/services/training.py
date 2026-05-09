"""
培训材料生成服务
生成PPT大纲和PPTX文件
"""

import json
import re
import uuid
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional

from backend.config import OUTPUT_DIR, get_kb_wiki_path
from backend.services.llm import llm_service


# PPT大纲生成Prompt
OUTLINE_PROMPT = """你是一个专业培训师。请基于以下知识库内容，生成一份培训PPT大纲。

## 培训要求

- 培训主题：{topic}
- 目标受众：{audience}
- 预计时长：{duration}分钟
- 幻灯片数量：{slide_count}页
- 内容侧重点：{focus_areas}

## 知识库内容

{content}

## 输出格式

请返回JSON格式的PPT大纲：

```json
{{
  "title": "培训标题",
  "chapters": [
    {{
      "title": "章节标题",
      "pages": 3,
      "points": [
        "要点1",
        "要点2",
        "要点3"
      ]
    }}
  ]
}}
```
"""


class TrainingService:
    """培训生成服务"""
    
    @staticmethod
    def _collect_knowledge_content(kb_ids: List[str]) -> str:
        """收集知识库内容"""
        content = []
        
        for kb_id in kb_ids:
            wiki_path = get_kb_wiki_path(kb_id)
            if not wiki_path.exists():
                continue
            
            for file_path in wiki_path.glob("*.md"):
                if file_path.name in ["index.md", "log.md"]:
                    continue
                
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
                
                # 只取前2000字
                content.append(f"## {file_path.stem}\n\n{text[:2000]}\n")
        
        return "\n\n".join(content)
    
    @staticmethod
    def _collect_document_content(doc_ids: List[str], kb_id: str) -> str:
        """收集文档内容"""
        from backend.services.document import doc_service
        
        # 这里简化处理，实际应该从raw/读取文档
        content = []
        
        # 读取知识库的Wiki页面作为文档内容
        wiki_path = get_kb_wiki_path(kb_id)
        if wiki_path.exists():
            for file_path in wiki_path.glob("*.md"):
                if file_path.name in ["index.md", "log.md"]:
                    continue
                
                with open(file_path, "r", encoding="utf-8") as f:
                    text = f.read()
                
                content.append(f"## {file_path.stem}\n\n{text[:2000]}\n")
        
        return "\n\n".join(content)
    
    @staticmethod
    async def generate_outline(
        source_type: str,
        source_ids: List[str],
        topic: str,
        audience: str,
        duration: int,
        slide_count: int,
        focus_areas: List[str],
        model_id: Optional[str] = None
    ) -> Dict:
        """
        生成PPT大纲
        
        Args:
            source_type: 来源类型 (knowledge_base, document)
            source_ids: 来源ID列表
            topic: 培训主题
            audience: 目标受众
            duration: 预计时长
            slide_count: 幻灯片数量
            focus_areas: 内容侧重点
            model_id: 模型ID
        
        Returns:
            PPT大纲
        """
        # 收集内容
        if source_type == "knowledge_base":
            content = TrainingService._collect_knowledge_content(source_ids)
        else:
            content = TrainingService._collect_document_content(source_ids, source_ids[0] if source_ids else "")
        
        if not content.strip():
            return {
                "title": topic,
                "chapters": [
                    {
                        "title": "第一章：概述",
                        "pages": 3,
                        "points": ["暂无内容，请先上传文档或创建知识库"]
                    }
                ]
            }
        
        # 构建Prompt
        prompt = OUTLINE_PROMPT.format(
            topic=topic,
            audience=audience,
            duration=duration,
            slide_count=slide_count,
            focus_areas=", ".join(focus_areas),
            content=content[:10000]  # 限制长度
        )
        
        messages = [
            {"role": "system", "content": "你是一个专业的企业安全培训师，擅长生成结构化的培训大纲。"},
            {"role": "user", "content": prompt}
        ]
        
        try:
            response = await llm_service.chat_sync(messages, model_id=model_id, temperature=0.7)
            
            # 提取JSON
            json_match = re.search(r'```json\s*(.*?)\s*```', response, re.DOTALL)
            if json_match:
                outline = json.loads(json_match.group(1))
            else:
                # 尝试直接解析JSON
                outline = json.loads(response)
            
            return outline
            
        except Exception as e:
            # 返回默认大纲
            return {
                "title": topic,
                "chapters": [
                    {
                        "title": "第一章：概述",
                        "pages": 3,
                        "points": [f"生成大纲时出错: {str(e)}", "请重试或调整参数"]
                    }
                ]
            }
    
    @staticmethod
    async def generate_ppt(
        outline: Dict,
        topic: str,
        audience: str,
        template: str = "default",
        model_id: Optional[str] = None
    ) -> str:
        """
        生成PPTX文件
        
        Args:
            outline: PPT大纲
            topic: 培训主题
            audience: 目标受众
            template: 模板名称
            model_id: 模型ID
        
        Returns:
            生成的PPTX文件路径
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            
            # 创建PPT
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # 标题页
            slide_layout = prs.slide_layouts[0]  # 标题页
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = outline.get("title", topic)
            subtitle.text = f"目标受众：{audience}\n基于安牛知识库自动生成"
            
            # 内容页
            for chapter in outline.get("chapters", []):
                # 章节标题页
                slide_layout = prs.slide_layouts[1]  # 标题和内容
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                body = slide.placeholders[1]
                
                title.text = chapter.get("title", "")
                
                tf = body.text_frame
                for point in chapter.get("points", []):
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = Pt(18)
            
            # 结束页
            slide_layout = prs.slide_layouts[6]  # 空白页
            slide = prs.slides.add_slide(slide_layout)
            
            # 添加感谢文字
            left = Inches(4)
            top = Inches(3)
            width = Inches(5)
            height = Inches(1.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.text = "谢谢！\nQuestions & Discussion"
            p = tf.paragraphs[0]
            p.font.size = Pt(36)
            p.alignment = PP_ALIGN.CENTER
            
            # 保存文件
            output_dir = OUTPUT_DIR
            output_dir.mkdir(parents=True, exist_ok=True)
            
            filename = f"{topic}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            output_path = output_dir / filename
            
            prs.save(str(output_path))
            
            return str(output_path)
            
        except Exception as e:
            raise Exception(f"PPT生成失败: {str(e)}")


# 服务实例
training_service = TrainingService()
