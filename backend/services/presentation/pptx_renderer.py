"""使用 python-pptx 渲染 PPTX。"""

from __future__ import annotations

import asyncio
import re
import textwrap
from threading import Event
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from backend.models import PresentationSpec, SlideSpec
from backend.services.presentation.project_store import get_job_paths, update_job_progress
from backend.services.presentation.safety_templates import SafetyTemplate

ACCENT_PALETTE = ["#0EA5E9", "#2563EB", "#7C3AED", "#F97316", "#10B981"]
BADGE_KEYWORDS = [
    "封面",
    "导入",
    "背景",
    "现状",
    "风险",
    "原因",
    "流程",
    "措施",
    "案例",
    "清单",
    "总结",
    "行动",
    "要点",
]


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _run_font(run, *, cn: str, en: str, size: int, color: str, bold: bool = False):
    run.font.name = cn
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _set_text(
    shape,
    text: str,
    *,
    cn: str,
    en: str,
    size: int,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            _run_font(run, cn=cn, en=en, size=size, color=color, bold=bold)
    return shape


def _set_paragraphs(
    shape,
    paragraphs: list[str],
    *,
    cn: str,
    en: str,
    size: int,
    color: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    spacing_after: int = 6,
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = valign
    cleaned_paragraphs = [re.sub(r"\s+", " ", str(p)).strip() for p in paragraphs if str(p).strip()]
    if not cleaned_paragraphs:
        cleaned_paragraphs = [""]
    for idx, paragraph_text in enumerate(cleaned_paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.text = paragraph_text
        p.space_after = Pt(spacing_after if idx < len(cleaned_paragraphs) - 1 else 0)
        for run in p.runs:
            _run_font(run, cn=cn, en=en, size=size, color=color, bold=bold)
    return shape


def _textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    cn,
    en,
    size,
    color,
    bold=False,
    align=PP_ALIGN.LEFT,
    fill: str | None = None,
    line: str | None = None,
    radius: bool = False,
    valign=MSO_ANCHOR.TOP,
):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(fill)
    else:
        box.fill.background()
    if line:
        box.line.color.rgb = _rgb(line)
    else:
        box.line.fill.background()
    return _set_text(box, text, cn=cn, en=en, size=size, color=color, bold=bold, align=align, valign=valign)


def _bg(slide, color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _safe_filename_fragment(value: str) -> str:
    fragment = "".join(ch for ch in value if ch.isalnum() or ch in {"-", "_", " "}).strip()
    fragment = re.sub(r"\s+", "_", fragment)
    return fragment[:80].strip("_")


def _card_badge_text(title: str, idx: int) -> str:
    cleaned = re.sub(r"\s+", "", str(title or ""))
    for keyword in BADGE_KEYWORDS:
        if keyword in cleaned:
            return keyword
    if cleaned:
        return cleaned[:2]
    return f"点{idx}"


def _wrap_display_text(text: str, *, line_width: int = 22, max_lines: int = 3) -> str:
    cleaned = re.sub(r"\s+", " ", str(text or "")).strip()
    cleaned = cleaned.replace("…", "").replace("...", "")
    if not cleaned:
        return ""
    lines = textwrap.wrap(cleaned, width=line_width, break_long_words=True, break_on_hyphens=False)
    if not lines:
        return cleaned
    return "\n".join(lines)


def _draw_frame(slide, tmpl: SafetyTemplate):
    _bg(slide, tmpl.theme_colors["bg"])
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32))
    header.fill.solid()
    header.fill.fore_color.rgb = _rgb(tmpl.theme_colors["primary"])
    header.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0.32), Inches(13.333), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(tmpl.theme_colors["accent"])
    accent.line.fill.background()


def _badge(slide, text: str, left: float, top: float, width: float, fill: str, fg: str, tmpl: SafetyTemplate):
    return _textbox(
        slide,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.34),
        text,
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=10,
        color=fg,
        bold=True,
        align=PP_ALIGN.CENTER,
        fill=fill,
        line=fill,
        radius=True,
        valign=MSO_ANCHOR.MIDDLE,
    )


def _bullet_card(
    slide,
    tmpl: SafetyTemplate,
    title: str,
    description: str,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    accent: str,
    badge_text: str,
):
    card = _textbox(
        slide,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        "",
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=12,
        color=tmpl.theme_colors["body"],
        fill="#FFFFFF",
        line=tmpl.theme_colors["border"],
        radius=True,
    )
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(0.12), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent)
    bar.line.fill.background()

    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + 0.24), Inches(top + 0.18), Inches(0.36), Inches(0.36))
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(accent)
    badge.line.fill.background()
    _set_text(
        badge,
        badge_text,
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=10,
        color="#FFFFFF",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    title_box = slide.shapes.add_textbox(Inches(left + 0.72), Inches(top + 0.12), Inches(width - 1.0), Inches(0.28))
    _set_text(
        title_box,
        title,
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=max(16, tmpl.body_size),
        color=tmpl.theme_colors["title"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    desc_box = slide.shapes.add_textbox(Inches(left + 0.72), Inches(top + 0.42), Inches(width - 1.0), Inches(height - 0.5))
    desc_box.text_frame.auto_size = None
    wrap_width = 30 if height <= 1.1 else 34
    _set_text(
        desc_box,
        _wrap_display_text(description, line_width=wrap_width),
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=max(14, tmpl.body_size - 1),
        color=tmpl.theme_colors["body"],
        valign=MSO_ANCHOR.TOP,
    )
    return card


def _render_text_page(slide, tmpl: SafetyTemplate, ss: SlideSpec):
    _draw_frame(slide, tmpl)
    title_box = slide.shapes.add_textbox(Inches(0.82), Inches(0.56), Inches(10.7), Inches(0.58))
    _set_text(
        title_box,
        ss.title,
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=tmpl.title_size,
        color=tmpl.theme_colors["title"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if ss.slide_type and ss.slide_type != "content":
        _badge(slide, ss.slide_type.replace("_", " "), 10.98, 0.6, 1.42, "#EEF2FF", tmpl.theme_colors["primary"], tmpl)

    subtitle = (ss.subtitle or ss.key_message or "").strip()
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.98), Inches(1.42), Inches(11.2), Inches(0.42))
        _set_text(
            subtitle_box,
            subtitle,
            cn=tmpl.font_family_cn,
            en=tmpl.font_family_en,
            size=max(16, tmpl.body_size + 1),
            color=tmpl.theme_colors["primary"],
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )

    body_paragraphs = [p for p in (ss.body_paragraphs or ss.bullets) if str(p).strip()]
    if not body_paragraphs:
        body_paragraphs = ["内容待补充"]

    body_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.82),
        Inches(1.92),
        Inches(11.72),
        Inches(4.75),
    )
    body_card.fill.solid()
    body_card.fill.fore_color.rgb = _rgb("#FFFFFF")
    body_card.line.color.rgb = _rgb(tmpl.theme_colors["border"])
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(1.92), Inches(0.12), Inches(4.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(ACCENT_PALETTE[0])
    bar.line.fill.background()

    body_box = slide.shapes.add_textbox(Inches(1.1), Inches(2.12), Inches(11.08), Inches(4.3))
    body_box.text_frame.auto_size = None
    _set_paragraphs(
        body_box,
        body_paragraphs,
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=max(14, tmpl.body_size - 1),
        color=tmpl.theme_colors["body"],
        bold=False,
        valign=MSO_ANCHOR.TOP,
        spacing_after=8,
    )
    return body_card


def _footer(slide, slide_no: int, tmpl: SafetyTemplate):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(12.08), Inches(7.02), Inches(0.62), Inches(0.28))
    pill.fill.solid()
    pill.fill.fore_color.rgb = _rgb("#F8FAFC")
    pill.line.color.rgb = _rgb(tmpl.theme_colors["border"])
    _set_text(
        pill,
        str(slide_no),
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=10,
        color=tmpl.theme_colors["body"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def _cover_labels(slide, tmpl: SafetyTemplate, ss: SlideSpec):
    labels = [item for item in ss.bullets if str(item).strip()][:3]
    if not labels:
        labels = [ss.title]
    card_top = 2.72
    card_height = 0.58
    for idx, raw_text in enumerate(labels[:3]):
        if "：" in raw_text:
            label, value = raw_text.split("：", 1)
        elif ":" in raw_text:
            label, value = raw_text.split(":", 1)
        else:
            label, value = "材料主题", raw_text
        top = card_top + idx * 0.74
        _textbox(
            slide,
            Inches(0.9),
            Inches(top),
            Inches(10.7),
            Inches(card_height),
            "",
            cn=tmpl.font_family_cn,
            en=tmpl.font_family_en,
            size=12,
            color=tmpl.theme_colors["body"],
            fill="#FFFFFF",
            line=tmpl.theme_colors["border"],
            radius=True,
        )
        _badge(slide, label.strip()[:8], 1.08, top + 0.12, 1.2, tmpl.theme_colors["primary"], "#FFFFFF", tmpl)
        value_box = slide.shapes.add_textbox(Inches(2.45), Inches(top + 0.09), Inches(8.75), Inches(0.28))
        _set_text(
            value_box,
            value.strip() or label.strip(),
            cn=tmpl.font_family_cn,
            en=tmpl.font_family_en,
            size=15,
            color=tmpl.theme_colors["title"],
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )


def _render_cover(slide, spec: PresentationSpec, tmpl: SafetyTemplate, ss: SlideSpec):
    _draw_frame(slide, tmpl)
    hero_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.76), Inches(0.78), Inches(11.82), Inches(1.78))
    hero_bg.fill.solid()
    hero_bg.fill.fore_color.rgb = _rgb("#FFFFFF")
    hero_bg.line.color.rgb = _rgb(tmpl.theme_colors["border"])
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.95), Inches(11.0), Inches(0.74))
    _set_text(
        title_box,
        spec.title,
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=tmpl.title_size + 8,
        color=tmpl.theme_colors["title"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    _cover_labels(slide, tmpl, ss)
    _footer(slide, 1, tmpl)


def _render_content(slide, tmpl: SafetyTemplate, ss: SlideSpec):
    if ss.visual_type == "text":
        _render_text_page(slide, tmpl, ss)
        _footer(slide, ss.slide_no, tmpl)
        return

    _draw_frame(slide, tmpl)
    title_box = slide.shapes.add_textbox(Inches(0.82), Inches(0.56), Inches(10.6), Inches(0.58))
    _set_text(
        title_box,
        ss.title,
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=tmpl.title_size,
        color=tmpl.theme_colors["title"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if ss.slide_type and ss.slide_type != "content":
        _badge(slide, ss.slide_type.replace("_", " "), 10.98, 0.6, 1.42, "#EEF2FF", tmpl.theme_colors["primary"], tmpl)

    bullets = [b for b in ss.bullets if str(b).strip()]
    if not bullets:
        bullets = ["内容待补充"]
    top = 1.46
    available_height = 5.95
    max_cards = min(len(bullets), 4)
    if max_cards <= 0:
        max_cards = 1
    card_height = min(1.68, max(1.2, (available_height - 0.12 * (max_cards - 1)) / max_cards))
    gap = 0.12
    for idx in range(max_cards):
        bullet = bullets[idx]
        if "：" in bullet:
            point_title, point_desc = bullet.split("：", 1)
        elif ":" in bullet:
            point_title, point_desc = bullet.split(":", 1)
        else:
            point_title, point_desc = bullet, ""
        point_desc = point_desc.strip()
        if point_desc and len(re.sub(r"\s+", "", point_desc)) < 10:
            point_desc = f"{point_desc}，用于现场执行和复核。"
        elif not point_desc:
            point_desc = "用于现场执行、检查和复盘。"
        if len(re.sub(r"\s+", "", point_desc)) < 18:
            point_desc = f"{point_desc}，并补充现场操作、检查方式和常见误区。"
        _bullet_card(
            slide,
            tmpl,
            point_title.strip()[:24],
            point_desc,
            left=0.82,
            top=top + idx * (card_height + gap),
            width=11.72,
            height=card_height,
            accent=ACCENT_PALETTE[idx % len(ACCENT_PALETTE)],
            badge_text=_card_badge_text(point_title.strip() or bullet, idx + 1),
        )

    if ss.slide_type == "quiz" and ss.notes:
        note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.4), Inches(6.3), Inches(4.14), Inches(0.5))
        note.fill.solid()
        note.fill.fore_color.rgb = _rgb("#FFFBEB")
        note.line.color.rgb = _rgb("#F59E0B")
        _set_text(
            note,
            f"答案提示：{ss.notes}",
            cn=tmpl.font_family_cn,
            en=tmpl.font_family_en,
            size=11,
            color="#92400E",
            valign=MSO_ANCHOR.MIDDLE,
        )

    _footer(slide, ss.slide_no, tmpl)


def _render_summary(slide, tmpl: SafetyTemplate, ss: SlideSpec):
    _draw_frame(slide, tmpl)
    title_box = slide.shapes.add_textbox(Inches(0.82), Inches(0.58), Inches(10.8), Inches(0.58))
    _set_text(
        title_box,
        ss.title,
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=tmpl.title_size,
        color=tmpl.theme_colors["title"],
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    bullets = [b for b in ss.bullets if str(b).strip()] or ["复盘关键风险点", "确认岗位职责和上报路径", "形成整改清单并闭环"]
    count = min(len(bullets), 3)
    available_height = 3.4
    card_height = min(1.25, max(1.0, (available_height - 0.12 * (count - 1)) / max(1, count)))
    for idx, bullet in enumerate(bullets[:3]):
        desc = "用于闭环整改、责任落实和复盘检查。"
        if "：" in bullet:
            _, maybe_desc = bullet.split("：", 1)
            if maybe_desc.strip():
                desc = f"{maybe_desc.strip()}，并纳入整改跟踪。"
        desc = desc if len(re.sub(r"\s+", "", desc)) >= 18 else f"{desc}建议按岗位逐项落实。"
        _bullet_card(
            slide,
            tmpl,
            bullet[:24],
            desc,
            left=0.82,
            top=1.58 + idx * (card_height + 0.12),
            width=11.72,
            height=card_height,
            accent=ACCENT_PALETTE[(idx + 1) % len(ACCENT_PALETTE)],
            badge_text=_card_badge_text(bullet, idx + 1),
        )
    closing = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(5.2), Inches(11.72), Inches(0.82))
    closing.fill.solid()
    closing.fill.fore_color.rgb = _rgb("#ECFDF5")
    closing.line.color.rgb = _rgb("#86EFAC")
    _set_text(
        closing,
        "建议：将本页内容转化为岗位清单、责任人和完成时限，直接进入闭环管理。",
        cn=tmpl.font_family_cn,
        en=tmpl.font_family_en,
        size=13,
        color="#166534",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    _footer(slide, ss.slide_no, tmpl)


def render_presentation(
    spec: PresentationSpec,
    template: SafetyTemplate,
    job_id: str,
    kb_name: str | None = None,
    filename_source: str | None = None,
    cancel_event: Event | None = None,
) -> dict[str, Any]:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(spec.slides)
    for idx, ss in enumerate(spec.slides, start=1):
        if cancel_event is not None and cancel_event.is_set():
            raise asyncio.CancelledError()

        update_job_progress(job_id, f"正在生成第 {idx}/{total} 页 PPT...")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if ss.slide_type == "cover":
            _render_cover(slide, spec, template, ss)
        elif ss.slide_type == "summary":
            _render_summary(slide, template, ss)
        else:
            _render_content(slide, template, ss)

    if cancel_event is not None and cancel_event.is_set():
        raise asyncio.CancelledError()

    paths = get_job_paths(job_id)
    paths.pptx_dir.mkdir(parents=True, exist_ok=True)
    source_name = filename_source or spec.title or spec.topic or job_id or "training"
    filename = f"{_safe_filename_fragment(source_name) or 'training'}.pptx"
    pptx_path = paths.pptx_dir / filename
    prs.save(str(pptx_path))
    return {
        "filename": filename,
        "pptx_path": str(pptx_path),
        "download_url": f"/api/training/download/{filename}",
    }
