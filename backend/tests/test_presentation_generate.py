"""PPT 生成全链路测试。"""

from __future__ import annotations

import asyncio
import os
import re
from threading import Event

import pytest
from pptx import Presentation

import backend.config as config_module
from backend.services.presentation.content_pack import build_content_pack
from backend.services.presentation.outline_builder import _build_llm_prompt, generate_outline
from backend.services.presentation.slide_planner import plan_slides
from backend.services.presentation.quality_check import check_presentation, repair_presentation
from backend.services.presentation.pptx_renderer import render_presentation
from backend.services.presentation.safety_templates import get_template
from backend.services.presentation.project_store import get_job_paths, resolve_download_path
from backend.models import PresentationSpec, SlideSpec, TrainingOutline


class PromptReq:
    sources = [{"type": "prompt", "prompt": "帮我生成一份危险化学品仓库火灾应急处置培训 PPT"}]
    topic = "危险化学品仓库火灾应急处置培训"
    audience = "一线员工"
    duration_minutes = 30
    slide_count = 15
    style = "standard_training"
    focus_areas = ["应急处置", "报警流程"]
    include_quiz = True
    template_id = "standard_training"


def test_outline_prompt_does_not_include_style_field(isolated_training_env):
    pack = build_content_pack(PromptReq())
    prompt = _build_llm_prompt(pack, {"title": "材料标题", "audience": "公司管理层", "style": "management_briefing", "slide_count": 6})

    assert "培训风格" not in prompt
    assert "管理层简报" not in prompt
    assert "一线班组" not in prompt
    assert "公司管理层" in prompt


def test_outline_to_pptx_full_chain(isolated_training_env, monkeypatch):
    monkeypatch.setitem(config_module.config["models"], "model_roles", {"ppt_gen": "deepseek-v4-flash"})
    monkeypatch.setitem(config_module.config["models"], "providers", [
        {"id": "test", "name": "Test", "base_url": "http://test", "api_key": "sk-test", "models": [{"id": "deepseek-v4-flash"}]}
    ])

    class MockLLM:
        async def chat_sync(self, messages, model_id=None, temperature=0.7, max_tokens=None):
            return '{"slides":[{"title":"危化品仓库风险","sections":[{"subtitle":"主要风险"}]},{"title":"应急处置流程","sections":[{"subtitle":"报警"},{"subtitle":"处置"}]},{"title":"总结","sections":[{"subtitle":"行动项"}]}]}'

    req = PromptReq()
    pack = build_content_pack(req)
    outline = asyncio.run(generate_outline(pack, req, llm_client=MockLLM()))
    spec = plan_slides(outline, pack, req)
    report = check_presentation(spec, pack, req)
    render_info = render_presentation(spec, get_template("standard_training"), "job-full-1")

    assert report.summary
    assert render_info["pptx_path"].endswith(".pptx")
    assert render_info["filename"] == f"{spec.title}.pptx"
    assert get_job_paths("job-full-1").pptx_dir.joinpath(render_info["filename"]).exists()


def test_render_presentation_uses_single_text_copy_and_no_brand_footer(isolated_training_env):
    spec = PresentationSpec(
        id="spec-cover",
        title="示例培训标题",
        topic="示例培训主题",
        audience="一线员工",
        duration_minutes=30,
        style="standard_training",
        template_id="standard_training",
        slides=[
            SlideSpec(
                id="slide-cover",
                slide_no=1,
                slide_type="cover",
                title="示例培训标题",
                bullets=["汇报时间：2026年5月21日", "汇报人：刘明超", "汇报对象：公司管理层"],
                source_refs=[],
                notes="",
                visual_type="cards",
                safety_level="normal",
            )
        ],
        quality_warnings=[],
    )
    render_info = render_presentation(spec, get_template("standard_training"), "job-text-check")
    prs = Presentation(render_info["pptx_path"])
    slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))

    assert slide_text.count("示例培训标题") == 1
    assert "AI 生成培训材料" not in slide_text
    assert "安牛" not in slide_text
    assert "..." not in slide_text
    assert "…" not in slide_text
    assert "2026年5月21日" in slide_text
    assert "刘明超" in slide_text
    assert "公司管理层" in slide_text


def test_render_presentation_does_not_show_core_message_block(isolated_training_env):
    spec = PresentationSpec(
        id="spec-core-message",
        title="示例标题",
        topic="示例主题",
        audience="管理层",
        duration_minutes=30,
        style="standard_training",
        template_id="standard_training",
        slides=[
            SlideSpec(
                id="slide-content",
                slide_no=1,
                slide_type="content",
                title="页面标题",
                bullets=["要点一：这里是一段更完整的说明文字，用于检查页面上不会再出现标题下注释。"],
                source_refs=[],
                notes="这是内部备注，不应展示为正文注释。",
                visual_type="cards",
                safety_level="normal",
            )
        ],
        quality_warnings=[],
    )

    render_info = render_presentation(spec, get_template("standard_training"), "job-core-check")
    prs = Presentation(render_info["pptx_path"])
    slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))

    assert "核心说明" not in slide_text
    assert "这是内部备注" not in slide_text
    assert "这里是一段更完整的说明文字" in slide_text


def test_render_presentation_renders_text_page_with_subtitle_and_body(isolated_training_env):
    spec = PresentationSpec(
        id="spec-text-page",
        title="文本页测试",
        topic="文本页测试",
        audience="管理层",
        duration_minutes=30,
        style="standard_training",
        template_id="standard_training",
        slides=[
            SlideSpec(
                id="slide-text",
                slide_no=1,
                slide_type="content",
                title="现场处置方案清单与适用场景",
                subtitle="把重点从提纲改成可直接朗读的正文",
                bullets=[],
                body_paragraphs=[
                    "供受油作业事故要先控制风险源，再明确人员疏散、现场隔离和信息上报顺序，避免处置动作混乱。",
                    "针对不同作业场景要补充班组分工和检查方式，确保每个人都知道自己的动作、时点和复核要求。",
                ],
                source_refs=[],
                notes="",
                visual_type="text",
                safety_level="normal",
            )
        ],
        quality_warnings=[],
    )

    render_info = render_presentation(spec, get_template("standard_training"), "job-text-page")
    prs = Presentation(render_info["pptx_path"])
    slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))

    assert "把重点从提纲改成可直接朗读的正文" in slide_text
    assert "供受油作业事故要先控制风险源" in slide_text
    assert "信息上报顺序" in slide_text


def test_render_presentation_keeps_long_bullet_text(isolated_training_env):
    spec = PresentationSpec(
        id="spec-long-bullet",
        title="长正文测试",
        topic="长正文测试",
        audience="管理层",
        duration_minutes=30,
        style="standard_training",
        template_id="standard_training",
        slides=[
            SlideSpec(
                id="slide-long",
                slide_no=1,
                slide_type="content",
                title="正文页",
                bullets=[
                    "要点一：这是一段很长的正文说明文字，用来验证 PPT 渲染时不会把后面的内容省略掉，而是完整写入并正确换行展示。",
                ],
                source_refs=[],
                notes="",
                visual_type="cards",
                safety_level="normal",
            )
        ],
        quality_warnings=[],
    )

    render_info = render_presentation(spec, get_template("standard_training"), "job-long-bullet")
    prs = Presentation(render_info["pptx_path"])
    slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
    normalized = re.sub(r"\s+", "", slide_text)

    assert "..." not in slide_text
    assert "…" not in slide_text
    assert "完整写入并正确换行展示" in normalized


def test_repair_presentation_preserves_long_bullet_text(isolated_training_env):
    spec = PresentationSpec(
        id="spec-repair-long-bullet",
        title="长正文测试",
        topic="长正文测试",
        audience="管理层",
        duration_minutes=30,
        style="standard_training",
        template_id="standard_training",
        slides=[
            SlideSpec(
                id="slide-repair-long",
                slide_no=1,
                slide_type="content",
                title="正文页",
                bullets=[
                    "要点一：领导要带头学习预案，将应急管理融入日常工作和班组检查中，确保每个环节都可执行、可检查、可复盘。",
                ],
                source_refs=[],
                notes="",
                visual_type="cards",
                safety_level="normal",
            )
        ],
        quality_warnings=[],
    )
    pack = build_content_pack(PromptReq())

    repaired = repair_presentation(spec, pack, {})
    assert repaired.slides[0].bullets[0].endswith("可执行、可检查、可复盘。")
    assert "…" not in repaired.slides[0].bullets[0]

    render_info = render_presentation(repaired, get_template("standard_training"), "job-repair-long-bullet")
    prs = Presentation(render_info["pptx_path"])
    slide_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
    normalized = re.sub(r"\s+", "", slide_text)

    assert "领导要带头学习预案，将应急管理融入日常工作和班组检查中，确保每个环节都可执行、可检查、可复盘。" in normalized
    assert "…" not in slide_text


def test_wrapped_text_keeps_consistent_font_size(isolated_training_env):
    spec = PresentationSpec(
        id="spec-font-consistency",
        title="字号一致性测试",
        topic="字号一致性测试",
        audience="管理层",
        duration_minutes=30,
        style="standard_training",
        template_id="standard_training",
        slides=[
            SlideSpec(
                id="slide-font",
                slide_no=1,
                slide_type="content",
                title="正文页",
                bullets=[
                    "要点一：供受油作业事故的现场处置方案，并补充现场操作、检查方式和常见误区，同时结合班组演练、现场隔离、人员疏散、信息上报和后续复盘进行闭环管理。",
                ],
                source_refs=[],
                notes="",
                visual_type="cards",
                safety_level="normal",
            )
        ],
        quality_warnings=[],
    )

    render_info = render_presentation(spec, get_template("standard_training"), "job-font-consistency")
    prs = Presentation(render_info["pptx_path"])
    slide = prs.slides[0]
    wrapped_shape = next(shape for shape in slide.shapes if hasattr(shape, "text") and "供受油作业事故" in shape.text)
    runs = wrapped_shape.text_frame.paragraphs[0].runs

    assert len(runs) >= 2
    sizes = {run.font.size for run in runs}
    assert len(sizes) == 1
    assert None not in sizes


def test_render_presentation_filename_uses_explicit_source_title(isolated_training_env):
    spec = PresentationSpec(
        id="spec-name",
        title="模型修正标题",
        topic="模型修正主题",
        audience="管理层",
        duration_minutes=30,
        style="standard_training",
        template_id="standard_training",
        slides=[],
        quality_warnings=[],
    )

    render_info = render_presentation(
        spec,
        get_template("standard_training"),
        "job-name-check",
        filename_source="用户输入的材料标题",
    )

    assert render_info["filename"] == "用户输入的材料标题.pptx"


def test_resolve_download_path_prefers_latest_matching_file(isolated_training_env):
    old_job = get_job_paths("job-old")
    new_job = get_job_paths("job-new")
    old_job.pptx_dir.mkdir(parents=True, exist_ok=True)
    new_job.pptx_dir.mkdir(parents=True, exist_ok=True)

    old_file = old_job.pptx_dir / "training_deck.pptx"
    new_file = new_job.pptx_dir / "training_deck.pptx"
    old_file.write_bytes(b"old")
    new_file.write_bytes(b"new")

    # 让旧文件更老一些，模拟历史 job 堆积在输出目录里的真实情况。
    older_ts = old_file.stat().st_mtime - 120
    os.utime(old_file, (older_ts, older_ts))

    resolved = resolve_download_path("training_deck.pptx")
    assert resolved == new_file.resolve()


def test_render_presentation_honors_cancel_event(isolated_training_env):
    spec = PresentationSpec(
        id="spec-cancel", title="测试", topic="测试", audience="管理层",
        duration_minutes=30, style="standard_training", template_id="standard_training",
        slides=[], quality_warnings=[],
    )
    cancel_event = Event()
    cancel_event.set()

    with pytest.raises(asyncio.CancelledError):
        render_presentation(spec, get_template("standard_training"), "job-cancel-1", cancel_event=cancel_event)


def test_quality_check_flags_missing_source(isolated_training_env):
    spec = PresentationSpec(
        id="spec-test", title="测试", topic="测试", audience="管理层",
        duration_minutes=30, style="standard_training", template_id="standard_training",
        slides=[
            SlideSpec(id="slide-1", slide_no=1, slide_type="legal_requirement", title="制度职责要求",
                      subtitle="", key_message="", bullets=["岗位职责必须明确"], source_refs=[],
                      notes="", visual_type="table", safety_level="warning")
        ], quality_warnings=[],
    )
    pack = build_content_pack({"sources": [{"type": "prompt", "prompt": "制度职责"}], "topic": "测试", "audience": "管理层"})
    report = check_presentation(spec, pack, {"slide_count": 1})
    assert any(issue.code in {"missing_source", "legal_without_source"} for issue in report.issues)


def test_ppt_service_generate_cancels_properly(isolated_training_env, monkeypatch):
    from backend.services.training_ppt import training_ppt_service

    cancel_event = Event()
    cancel_event.set()

    monkeypatch.setattr(
        "backend.services.training_ppt.get_running_job_cancel_event",
        lambda job_id: cancel_event,
    )

    async def fake_build_outline(*args, **kwargs):
        return TrainingOutline(
            id="outline-cancel", title="测试取消", topic="测试取消",
            audience="一线员工", duration_minutes=30, style="standard_training",
            slides=[], sections=[], warnings=[],
        )

    async def fake_render(*args, **kwargs):
        raise AssertionError("render should not be reached after cancellation")

    monkeypatch.setattr("backend.services.training_ppt.render_presentation", fake_render)
    monkeypatch.setattr("backend.services.training_ppt.build_outline", fake_build_outline)

    async def run():
        payload = {
            "sources": [{"type": "prompt", "prompt": "应急处置培训"}],
            "title": "测试取消", "topic": "测试取消",
            "audience": "一线员工", "duration_minutes": 30,
            "slide_count": 6, "style": "standard_training",
            "template_id": "standard_training", "include_quiz": True,
        }
        with pytest.raises(asyncio.CancelledError):
            await training_ppt_service.generate_ppt(payload, job_id="job-cancel-render")

    asyncio.run(run())
